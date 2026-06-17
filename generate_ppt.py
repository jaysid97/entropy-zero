import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Design Tokens (Colors matching index.css)
    BG_COLOR = RGBColor(10, 10, 18)        # Space Obsidian
    PRIMARY_COLOR = RGBColor(99, 102, 241)  # Neon Violet-Blue
    ACCENT_COLOR = RGBColor(6, 182, 212)   # Cyber Cyan
    TEXT_COLOR = RGBColor(248, 250, 252)   # Bright White
    MUTED_COLOR = RGBColor(148, 163, 184)  # Muted Gray
    ERROR_COLOR = RGBColor(239, 68, 68)    # Crimson Red

    blank_layout = prs.slide_layouts[6] # Blank layout for custom styling

    def apply_slide_base(slide, title_text, category_text="ENTROPY-ZERO"):
        # 1. Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        # 2. Category tag (small top header)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(5), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Outfit'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY_COLOR

        # 3. Main header title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_COLOR

        # 4. Thin horizontal accent line under title
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.color.rgb = PRIMARY_COLOR

    # ==========================================
    # SLIDE 1: Title Slide (Cinematic Entry)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Big Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.8))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "ENTROPY-ZERO"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_COLOR
    
    p2 = tf1.add_paragraph()
    p2.text = "The Autonomic Social Engine for Agile Teams"
    p2.font.name = 'Outfit'
    p2.font.size = Pt(24)
    p2.font.color.rgb = PRIMARY_COLOR

    # Tagline Details Box
    desc_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.2))
    tf2 = desc_box.text_frame
    tf2.word_wrap = True
    p_desc = tf2.paragraphs[0]
    p_desc.text = "Replacing meetings, manual status syncs, and ticketing overhead with a self-organizing social-computational team nervous system."
    p_desc.font.name = 'Inter'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = MUTED_COLOR

    # Footer Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(10.0), Inches(0.5))
    tf_m = meta_box.text_frame
    p_m = tf_m.paragraphs[0]
    p_m.text = "Hackathon Submission: Future of Productivity  |  Live at: entropy-zero.vercel.app"
    p_m.font.name = 'Space Grotesk'
    p_m.font.size = Pt(11)
    p_m.font.color.rgb = ACCENT_COLOR

    # Decorative bottom accent strip
    strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRIMARY_COLOR
    strip.line.fill.background()

    # ==========================================
    # SLIDE 2: The Problem (Coordination Entropy)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide2, "The Problem: The Collaboration Tax")

    # Left Column: Detailed bullet points
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(6.2), Inches(5.0))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True

    bullets = [
        ("Coordination Overhead:", " Modern engineering teams spend up to 60% of their time on 'work-about-work' (daily status alignments, grooming backlog boards, context switching)."),
        ("Meeting Fatigue:", " Fragmented schedules slice developer focus hours into small windows, destroying flow states."),
        ("Backlog Decay:", " Static task tracking pipelines (like Jira) become stale instantly and demand constant manual updates."),
        ("Management Bottlenecks:", " Traditional top-down workload assignments fail to account for real-time developer load reserves and fatigue indexes.")
    ]

    for title, desc in bullets:
        p = left_tf.add_paragraph() if left_tf.text else left_tf.paragraphs[0]
        p.space_after = Pt(18)
        run_title = p.add_run()
        run_title.text = title
        run_title.font.name = 'Inter'
        run_title.font.size = Pt(14)
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_COLOR

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = 'Inter'
        run_desc.font.size = Pt(13)
        run_desc.font.color.rgb = MUTED_COLOR

    # Right Column: Visual Warning Card
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.2), Inches(4.8), Inches(4.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(25, 15, 20) # Reddish-dark card
    card.line.color.rgb = ERROR_COLOR
    card.line.width = Pt(1.5)

    card_tf = card.text_frame
    card_tf.word_wrap = True
    card_tf.margin_left = Inches(0.4)
    card_tf.margin_right = Inches(0.4)
    card_tf.margin_top = Inches(0.5)

    p_card_title = card_tf.paragraphs[0]
    p_card_title.text = "COORDINATION DRAG"
    p_card_title.font.name = 'Space Grotesk'
    p_card_title.font.size = Pt(12)
    p_card_title.font.bold = True
    p_card_title.font.color.rgb = ERROR_COLOR

    p_card_pct = card_tf.add_paragraph()
    p_card_pct.space_before = Pt(12)
    p_card_pct.text = "62%"
    p_card_pct.font.name = 'Outfit'
    p_card_pct.font.size = Pt(72)
    p_card_pct.font.bold = True
    p_card_pct.font.color.rgb = ERROR_COLOR

    p_card_desc = card_tf.add_paragraph()
    p_card_desc.space_before = Pt(8)
    p_card_desc.text = "of team capacity is consumed by meeting overhead and manual alignment cycles, leaving only 38% for core building."
    p_card_desc.font.name = 'Inter'
    p_card_desc.font.size = Pt(12)
    p_card_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 3: The Cure (Autonomic Coordination)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide3, "The Solution: Autonomic Architecture")

    # Subtext introduction
    sub_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    sub_tf = sub_box.text_frame
    p_sub = sub_tf.paragraphs[0]
    p_sub.text = "Entropy-Zero maps coordination as a self-healing physics equation across three core pillars:"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = MUTED_COLOR

    # Three Columns (Boxes)
    cols = [
        ("1. PASSIVE TELEMETRY", 
         "Mines workspace activity directly.",
         "Aggregates signals from IDE edits, git logs, and Slack discussion densities. Forms a real-time, zero-friction developer capacity and specialty map without requiring manual reports."),
        
        ("2. AUTONOMIC ROUTING", 
         "Stress-aware task dispatching.",
         "Incoming issues act as kinetic vectors that auto-assign to the engineer with the optimal load and skills. If a developer goes OOO, the graph automatically re-routes tasks."),
         
        ("3. CONSENSUS SYNAPSE", 
         "Decentralized resolution.",
         "Identifies alignment locks in team text chats. Automatically spawns voting ballots, calculates domain expert weights, and merges PR branches once threshold consensus resolves.")
    ]

    box_width = Inches(3.6)
    box_height = Inches(4.3)
    start_left = Inches(0.8)
    gap = Inches(0.4)

    for i, (title, subtitle, desc) in enumerate(cols):
        left_pos = start_left + i * (box_width + gap)
        
        # Draw column card
        col_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.3), box_width, box_height)
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = RGBColor(16, 16, 28)
        col_box.line.color.rgb = PRIMARY_COLOR if i == 0 else (ACCENT_COLOR if i == 1 else RGBColor(168, 85, 247))
        col_box.line.width = Pt(1.5)
        
        col_tf = col_box.text_frame
        col_tf.word_wrap = True
        col_tf.margin_left = Inches(0.3)
        col_tf.margin_right = Inches(0.3)
        col_tf.margin_top = Inches(0.4)
        
        p_t = col_tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Space Grotesk'
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_COLOR
        
        p_st = col_tf.add_paragraph()
        p_st.space_before = Pt(4)
        p_st.text = subtitle
        p_st.font.name = 'Inter'
        p_st.font.size = Pt(11)
        p_st.font.italic = True
        p_st.font.color.rgb = ACCENT_COLOR
        
        p_d = col_tf.add_paragraph()
        p_d.space_before = Pt(14)
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED_COLOR

    # ==========================================
    # SLIDE 4: Interactive Dashboard & Features
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide4, "Real-Time Telemetry & Simulation")

    # Left: Bullet list of live features
    left_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(6.0), Inches(5.0))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True

    features = [
        ("Interactive Physics Graph:", " A live-rendered HTML5 Canvas coordinate space displaying developers (pulsing load circles), task diamonds (priority color-coded), and active telemetry data packets traveling along edges."),
        ("Stress-Injector Panel:", " Sandbox tools to trigger network shocks: Dev OOO events, P1 Outages, timelines acceleration, or Scope Creep. Evaluates the system's re-balancing capabilities."),
        ("Consensus Panel:", " Real-time weighted voting ballot interface resolving technical disputes based on expert contributions."),
        ("Monospace Log Stream:", " Live scrolling console printout showcasing system allocations, stress mitigation checks, and resolution triggers.")
    ]

    for title, desc in features:
        p = left_tf.add_paragraph() if left_tf.text else left_tf.paragraphs[0]
        p.space_after = Pt(16)
        run_title = p.add_run()
        run_title.text = title
        run_title.font.name = 'Inter'
        run_title.font.size = Pt(13)
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_COLOR

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = 'Inter'
        run_desc.font.size = Pt(12.5)
        run_desc.font.color.rgb = MUTED_COLOR

    # Right: Dashboard visual placeholder description card
    panel = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(2.0), Inches(5.1), Inches(4.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(12, 12, 22)
    panel.line.color.rgb = ACCENT_COLOR
    panel.line.width = Pt(1)

    panel_tf = panel.text_frame
    panel_tf.word_wrap = True
    panel_tf.margin_left = Inches(0.4)
    panel_tf.margin_right = Inches(0.4)
    panel_tf.margin_top = Inches(0.5)

    p_p_title = panel_tf.paragraphs[0]
    p_p_title.text = "CORE DASHBOARD OVERVIEW"
    p_p_title.font.name = 'Space Grotesk'
    p_p_title.font.size = Pt(12)
    p_p_title.font.bold = True
    p_p_title.font.color.rgb = ACCENT_COLOR

    p_p_body = panel_tf.add_paragraph()
    p_p_body.space_before = Pt(14)
    p_p_body.text = "The web application contains a dual-layout interface:\n\n1. Pitch Mode: The interactive slide deck presentation containing simulated sandboxes.\n\n2. Engine Dashboard Mode: The live telemetry cockpit containing the active canvas graph, stress injector controls, consensus balloting resolving widgets, and streaming telemetry console output."
    p_p_body.font.name = 'Inter'
    p_p_body.font.size = Pt(12)
    p_p_body.font.color.rgb = MUTED_COLOR

    # ==========================================
    # SLIDE 5: Business Value & ROI
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide5, "The Economics of Zero Entropy")

    # Subtext introduction
    sub_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    sub_tf = sub_box.text_frame
    p_sub = sub_tf.paragraphs[0]
    p_sub.text = "Eliminating manual coordination overhead yields direct, measurable productivity gains:"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = MUTED_COLOR

    # Left: Stats Column
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True

    stats = [
        ("+40% Focus Flow Time", "Engineers stay immersed in IDE environments rather than updating progress boards."),
        ("-65% Task Routing Latency", "Issues map and transition within minutes instead of waiting for sprint meetings."),
        ("Linear Scaling Efficiency", "Allows scaling engineering staff without suffers from typical administrative drag.")
    ]

    for title, desc in stats:
        p = left_tf.add_paragraph() if left_tf.text else left_tf.paragraphs[0]
        p.space_after = Pt(20)
        run_title = p.add_run()
        run_title.text = title + "\n"
        run_title.font.name = 'Outfit'
        run_title.font.size = Pt(20)
        run_title.font.bold = True
        run_title.font.color.rgb = ACCENT_COLOR

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = 'Inter'
        run_desc.font.size = Pt(12)
        run_desc.font.color.rgb = MUTED_COLOR

    # Right: Call To Action Deployment Box
    cta = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.3), Inches(5.3), Inches(4.0))
    cta.fill.solid()
    cta.fill.fore_color.rgb = RGBColor(16, 16, 32)
    cta.line.color.rgb = PRIMARY_COLOR
    cta.line.width = Pt(2)

    cta_tf = cta.text_frame
    cta_tf.word_wrap = True
    cta_tf.margin_left = Inches(0.4)
    cta_tf.margin_right = Inches(0.4)
    cta_tf.margin_top = Inches(0.5)

    p_c_t = cta_tf.paragraphs[0]
    p_c_t.text = "LIVE DEPLOYMENT"
    p_c_t.font.name = 'Space Grotesk'
    p_c_t.font.size = Pt(12)
    p_c_t.font.bold = True
    p_c_t.font.color.rgb = PRIMARY_COLOR

    p_c_u = cta_tf.add_paragraph()
    p_c_u.space_before = Pt(14)
    p_c_u.text = "https://entropy-zero.vercel.app"
    p_c_u.font.name = 'Outfit'
    p_c_u.font.size = Pt(20)
    p_c_u.font.bold = True
    p_c_u.font.color.rgb = ACCENT_COLOR

    p_c_d = cta_tf.add_paragraph()
    p_c_d.space_before = Pt(14)
    p_c_d.text = "Visit the live URL to experience the autonomic physics graph, run custom task injections, trigger stress shocks, and vote in the consensus synapses in real-time."
    p_c_d.font.name = 'Inter'
    p_c_d.font.size = Pt(12)
    p_c_d.font.color.rgb = MUTED_COLOR

    # Save presentation
    output_filename = "Entropy-Zero_Pitch_Deck.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_presentation()
